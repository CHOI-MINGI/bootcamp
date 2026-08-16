"""강의 슬라이드(PPTX) 생성 기능.

역할 분담:
  - LLM   : 자료를 근거로 "무엇을 넣을지"(제목·내용·출처)를 JSON으로 정한다.
  - 이 파일 : 그 JSON을 받아 "어떻게 배치할지"를 처리해 실제 pptx 파일을 만든다.

그림은 새로 생성하지 않고, 근거가 된 PDF 페이지에서 도표·이미지를 그대로 가져온다.
자료에 없는 그림을 지어내지 않기 위해서다.
"""

import io
import json
import re

import fitz                                  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

from rag_module import (_search, _build_context, _build_sources,
                        _format_extra, _failed, get_llm)


# ============================================================
# 1. 색상 / 서체
# ============================================================
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x12, 0x19, 0x3F)
NAVY_MID = RGBColor(0x2A, 0x36, 0x70)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
ICE_LIGHT = RGBColor(0xEA, 0xF1, 0xFC)
AMBER = RGBColor(0xD9, 0x8A, 0x1F)
AMBER_LIGHT = RGBColor(0xFB, 0xF0, 0xDC)
GRAY = RGBColor(0x5A, 0x64, 0x78)
LINE = RGBColor(0xD8, 0xDE, 0xE9)
TEXT = RGBColor(0x1A, 0x21, 0x38)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF3, 0xF6, 0xFB)

FONT = "맑은 고딕"
MONO = "Consolas"

SLIDE_W = 13.333
SLIDE_H = 7.5

# 내용의 성격에 따라 배치를 달리한다. 전부 같은 모양이면 읽는 사람이 지친다.
SLIDE_TYPES = ("개념", "절차", "비교", "코드")


# ============================================================
# 2. LLM에게 슬라이드 구성을 받아오는 부분
# ============================================================
# 두 단계로 나눈다.
#   1) 아웃라인 — 어떤 제목으로 몇 장을 만들지 먼저 정한다.
#   2) 본문 — 제목마다 다시 검색해 그 슬라이드에 맞는 근거를 붙여 채운다.
#
# 한 번의 검색 결과로 전체를 만들면 슬라이드마다 필요한 근거가 다른데도
# 같은 재료를 나눠 쓰게 되어 뒤로 갈수록 내용이 얕아진다.
OUTLINE_TEMPLATE = """당신은 OO대학교의 강의자료 기획 AI입니다.

[할 일]
주제 '{topic}'로 강의 슬라이드 {count}장을 만들려고 합니다.
각 슬라이드의 제목만 먼저 정하십시오. 내용은 아직 쓰지 마십시오.

[규칙]
1. 아래 [자료]에 실제로 있는 내용만 다루십시오.
2. 강의 흐름을 따르십시오. 개념 정의 → 원리 → 적용·절차 → 주의점 순이 자연스럽습니다.
3. 제목은 20자 이내의 명사구로 적으십시오. 학생이 나중에 목차로 찾을 수 있어야 합니다.
4. 제목끼리 내용이 겹치지 않게 하십시오.
5. JSON 배열만 출력하십시오. 형식은 다음과 같습니다.
   ["전류와 전압의 정의", "키르히호프 전류 법칙", "절점 해석법의 순서"]
{extra}
[자료]
{context}

[JSON 배열만 출력]"""


SLIDE_TEMPLATE = """당신은 OO대학교의 강의자료 작성 AI입니다.

[규칙]
1. 아래 [자료]에 있는 내용만 사용하십시오. 자료에 없는 내용은 절대 지어내지 마십시오.
2. 반드시 JSON 배열만 출력하십시오. 설명 문장, 인사말, 코드블록 표시를 붙이지 마십시오.
3. 각 슬라이드는 다음 형식을 따르십시오.
   {{"type": "개념", "title": "슬라이드 제목", "lead": "한 줄 요약",
     "bullets": ["문장1", "문장2"], "code": "", "diagram": null,
     "source": "파일명 p.페이지"}}
4. type은 내용의 성격에 따라 아래 넷 중 하나로 정하십시오.
   개념  무엇인지 설명하는 내용
   절차  순서가 있는 단계. bullets를 순서대로 씁니다
   비교  둘을 견주는 내용. bullets를 "A 항목 | B 항목" 형태로 짝지어 씁니다
   코드  코드나 태그가 핵심인 내용. code 항목에 원문을 넣습니다
5. lead는 그 슬라이드에서 학습자가 얻어야 할 핵심을 한 문장(35자 이내)으로 적습니다.
   bullets에 나올 문장을 그대로 옮기지 마십시오. 겹치면 같은 말을 두 번 하는 셈입니다.
6. bullets는 슬라이드당 3~4개, 각 문장은 45자 이내로 작성하십시오.
   서로 다른 내용을 담아야 하며, 표현만 바꾼 반복은 넣지 마십시오.
7. code는 type이 "코드"일 때만 채웁니다. 자료에 있는 코드를 그대로 옮기고,
   12줄을 넘지 않게 핵심 부분만 발췌하십시오. 그 외에는 빈 문자열로 두십시오.
8. 본문에서도 태그·속성·명령어는 원문 그대로 표기하십시오. 예: <script>, background-color
9. diagram은 적극적으로 채우십시오. 그림이 있는 슬라이드가 글만 있는 슬라이드보다 낫습니다.
   자료에서 아래 중 하나라도 읽히면 반드시 만드십시오.
     순서나 단계가 있다             → 흐름
     원인과 결과, 입력과 출력이 있다 → 흐름
     돌아오거나 반복되는 관계가 있다 → 순환
     상위 개념 아래 항목이 나뉜다    → 계층
   {{"kind": "흐름", "nodes": ["웹 클라이언트", "인터넷", "웹 서버"],
     "edges": [[0, 1, "요청"], [2, 1, "응답"]]}}
   nodes는 3~5개, 각 이름은 12자 이내로 짧게 적으십시오.
   edges의 [출발, 도착, "설명"]에서 번호는 nodes의 순서(0부터)입니다.
   설명이 필요 없으면 빈 문자열로 두고, 8자를 넘기지 마십시오.
   자료에 없는 관계를 지어내서는 안 됩니다. 다만 자료에 있는 관계를 그냥 넘기지도 마십시오.
   관계를 정말 찾을 수 없을 때만 null로 두십시오.
10. source에는 그 슬라이드의 근거가 된 출처를 '파일명 p.숫자' 형식으로 정확히 적으십시오.
11. 같은 type이 연속으로 두 번 나오지 않게 구성하십시오.
12. type이 "개념"인 슬라이드는 전체의 절반을 넘기지 마십시오.
13. 모든 내용은 한국어로 작성하십시오.

[예시]
아래는 다른 주제로 잘 작성된 슬라이드입니다. 형식과 밀도를 이 정도로 맞추십시오.
[
  {{"type": "개념", "title": "관계형 데이터베이스",
    "lead": "데이터를 표로 나누고 열쇠로 이어 붙인다",
    "bullets": ["행은 하나의 사례, 열은 하나의 속성을 뜻한다",
                "기본키는 행을 중복 없이 가리키는 열이다",
                "외래키는 다른 표의 기본키를 가리켜 두 표를 잇는다"],
    "code": "",
    "diagram": {{"kind": "계층", "nodes": ["데이터베이스", "테이블", "행", "열"],
                "edges": [[0, 1, "포함"], [1, 2, ""], [1, 3, ""]]}},
    "source": "DB개론.pdf p.12"}},
  {{"type": "절차", "title": "질의 처리 순서",
    "lead": "SQL 한 줄은 네 단계를 거쳐 결과가 된다",
    "bullets": ["구문 분석에서 문법 오류를 먼저 걸러낸다",
                "최적화기가 여러 실행 경로 중 비용이 낮은 것을 고른다",
                "실행기가 선택된 경로대로 데이터를 읽는다",
                "결과를 정해진 형식으로 묶어 돌려준다"],
    "code": "",
    "diagram": {{"kind": "흐름", "nodes": ["질의 입력", "구문 분석", "최적화", "실행", "결과 반환"],
                "edges": [[0, 1, ""], [1, 2, ""], [2, 3, "실행 계획"], [3, 4, ""]]}},
    "source": "DB개론.pdf p.31"}}
]

[요청]
주제 '{topic}'에 대한 강의 슬라이드 {count}장을 구성하십시오.
[자료]에 슬라이드 번호와 제목이 미리 정해져 있으면 그 제목을 그대로 쓰고,
각 슬라이드는 자기 번호에 배정된 근거를 우선 사용하십시오.
{extra}
[자료]
{context}

[JSON 배열만 출력]"""


DIAGRAM_KINDS = ("흐름", "순환", "계층")

MAX_NODES = 5
MAX_NODE_LABEL = 12
MAX_EDGE_LABEL = 8


def _clean_diagram(raw):
    """LLM이 준 도식 구조를 검사해 그릴 수 있는 형태로 다듬는다.

    노드가 너무 많거나 이름이 길면 도형 밖으로 글자가 넘친다.
    화살표가 없는 노드 목록은 도식이라 할 수 없으므로 버린다.
    """
    if not isinstance(raw, dict):
        return None

    kind = str(raw.get("kind", "")).strip()
    if kind not in DIAGRAM_KINDS:
        kind = "흐름"

    nodes = [str(n).strip()[:MAX_NODE_LABEL]
             for n in raw.get("nodes", []) if str(n).strip()]
    nodes = nodes[:MAX_NODES]

    if len(nodes) < 2:
        return None

    edges = []
    for e in raw.get("edges", []):
        if not isinstance(e, (list, tuple)) or len(e) < 2:
            continue
        try:
            a, b = int(e[0]), int(e[1])
        except (TypeError, ValueError):
            continue

        # 없는 노드를 가리키거나 자기 자신으로 도는 화살표는 버린다.
        if not (0 <= a < len(nodes)) or not (0 <= b < len(nodes)) or a == b:
            continue

        label = str(e[2]).strip()[:MAX_EDGE_LABEL] if len(e) > 2 else ""
        edges.append((a, b, label))

    if not edges:
        return None

    return {"kind": kind, "nodes": nodes, "edges": edges}


def _parse_json(raw):
    """LLM 응답에서 JSON 배열을 꺼낸다.

    지시를 했더라도 앞뒤에 설명이나 ```json 표시가 붙는 경우가 있어,
    가장 바깥쪽 대괄호 구간만 잘라내어 파싱한다.
    """
    text = raw.strip()
    text = re.sub(r"^```(?:json)?|```$", "", text, flags=re.MULTILINE).strip()

    start = text.find("[")
    end = text.rfind("]")
    if start == -1 or end == -1:
        raise ValueError("JSON 배열을 찾지 못했습니다.")

    data = json.loads(text[start:end + 1])

    cleaned = []
    for item in data:
        if not isinstance(item, dict):
            continue

        kind = str(item.get("type", "")).strip()
        if kind not in SLIDE_TYPES:
            kind = "개념"

        code = str(item.get("code", "") or "").strip()
        if code and kind != "코드":
            kind = "코드"          # 코드를 넣었으면 그 유형으로 다룬다

        cleaned.append({
            "type": kind,
            "title": str(item.get("title", "")).strip(),
            "lead": str(item.get("lead", "")).strip(),
            "bullets": [str(b).strip() for b in item.get("bullets", []) if str(b).strip()],
            "code": code,
            "diagram": _clean_diagram(item.get("diagram")),
            "source": str(item.get("source", "")).strip(),
        })

    if not cleaned:
        raise ValueError("슬라이드 항목이 비어 있습니다.")

    return cleaned


def _parse_titles(raw, count):
    """아웃라인 응답에서 제목 목록을 꺼낸다."""
    text = re.sub(r"^```(?:json)?|```$", "", raw.strip(), flags=re.MULTILINE).strip()

    start, end = text.find("["), text.rfind("]")
    if start == -1 or end == -1:
        raise ValueError("제목 배열을 찾지 못했습니다.")

    titles = []
    for item in json.loads(text[start:end + 1]):
        title = str(item).strip()
        if title and title not in titles:
            titles.append(title)

    if not titles:
        raise ValueError("제목이 비어 있습니다.")

    return titles[:count]


def _outline_context(docs):
    """아웃라인 단계용 자료 묶음.

    본문 단계와 달리 유사도 순이 아니라 '자료에 실린 순서'로 늘어놓는다.
    앞에서 뒤로 읽히면 문서가 어떤 순서로 전개되는지가 드러나므로,
    LLM이 강의 흐름에 맞는 목차를 잡기 쉬워진다.
    """
    def position(doc):
        meta = doc.metadata or {}
        return (str(meta.get("source", "")), meta.get("page", 0) or 0)

    return _build_context(sorted(docs, key=position))


def _slide_context(pairs):
    """슬라이드별로 배정된 근거를 번호와 제목을 붙여 늘어놓는다."""
    blocks = []
    for i, (title, docs) in enumerate(pairs, start=1):
        blocks.append(f"[슬라이드 {i}] 제목: {title}\n{_build_context(docs)}")
    return "\n\n".join(blocks)


def generate_slide_data(vectorstore, topic, count=5, role="교수자", courses=(),
                        k=6, extra=""):
    """자료를 근거로 슬라이드 구성(JSON)을 만든다.

    두 단계로 나눈다.
      1) 아웃라인 — 넓게 검색해 슬라이드 제목을 먼저 정한다.
      2) 본문 — 제목마다 다시 검색해 그 슬라이드에 맞는 근거를 붙여 채운다.

    한 번 검색한 결과를 여러 장이 나눠 쓰면 뒤로 갈수록 쓸 내용이 떨어진다.
    검색은 임베딩 계산이라 값이 싸므로, 제목 수만큼 다시 찾아도 부담이 적다.
    LLM 호출만 두 번으로 묶어 비용을 억제한다.

    extra : 교수자가 입력한 추가 지시사항 (선택)
    """
    # 아웃라인 단계는 문서 전체 구성을 봐야 하므로 더 넓게 가져온다.
    docs, best_score, blocked = _search(vectorstore, topic, role, courses, max(k, count * 2))
    if blocked:
        return blocked

    # ---------- 1단계: 제목 정하기 ----------
    outline_chain = (ChatPromptTemplate.from_template(OUTLINE_TEMPLATE)
                     | get_llm(0.2) | StrOutputParser())
    try:
        titles = _parse_titles(outline_chain.invoke({
            "topic": topic,
            "count": count,
            "extra": _format_extra(extra),
            "context": _outline_context(docs),
        }), count)
    except Exception:
        # 아웃라인이 실패해도 기능을 멈추지 않는다.
        # 제목 없이 예전 방식대로 한 번에 만들면 결과가 조금 얕아질 뿐이다.
        titles = []

    # ---------- 2단계: 제목별 근거 모으기 ----------
    if titles:
        pairs, seen = [], []
        for title in titles:
            # 제목만으로 검색하면 주제에서 벗어날 수 있어 주제를 함께 넣는다.
            found, _, sub_blocked = _search(vectorstore, f"{topic} {title}",
                                            role, courses, 4)
            # 가드레일에 걸리면 앞서 받아둔 자료로 대신한다.
            pairs.append((title, found if not sub_blocked and found else docs[:3]))
            seen.extend(pairs[-1][1])

        context = _slide_context(pairs)

        # 출처 목록은 실제로 쓰인 자료 전체를 기준으로 만든다.
        # 검색마다 새 객체가 만들어지므로 내용으로 중복을 걸러낸다.
        unique = {}
        for doc in seen:
            meta = doc.metadata or {}
            key = (meta.get("source"), meta.get("page"), doc.page_content[:40])
            unique.setdefault(key, doc)
        docs = list(unique.values())
    else:
        context = _build_context(docs)

    # ---------- 3단계: 본문 채우기 ----------
    prompt = ChatPromptTemplate.from_template(SLIDE_TEMPLATE)
    # 슬라이드는 JSON 형식을 지켜야 하므로 생성 기능 중 가장 낮은 값을 쓴다.
    chain = prompt | get_llm(0.2) | StrOutputParser()

    try:
        raw = chain.invoke({
            "topic": topic,
            "count": count,
            "extra": _format_extra(extra),
            "context": context,
        })
    except Exception as e:
        result = _failed(e)
        result["slides"] = None
        return result

    try:
        slides = _parse_json(raw)
    except Exception as e:
        # JSON 파싱에 실패해도 기능이 통째로 죽지 않도록 원문을 돌려준다.
        return {
            "answer": raw,
            "slides": None,
            "sources": _build_sources(docs),
            "blocked": False,
            "reason": f"JSON 파싱 실패: {e}",
            "best_score": best_score,
        }

    return {
        "answer": None,
        "slides": slides,
        "sources": _build_sources(docs),
        "blocked": False,
        "reason": None,
        "best_score": best_score,
    }


# ============================================================
# 3. 근거 페이지에서 그림 가져오기
# ============================================================
def _parse_source(source):
    """'회로이론.pdf p.12' 형태에서 파일명과 페이지 번호를 분리한다.

    도표를 꺼낼 수 있는 것은 PDF뿐이므로 다른 형식은 처리하지 않는다.
    """
    m = re.search(r"^(.*?\.pdf)\s*p\.\s*(\d+)", source, re.IGNORECASE)
    if not m:
        return None, None
    return m.group(1).strip(), int(m.group(2))


def _image_stats(image_bytes):
    """이미지의 (내용 비율, 흑백 여부)를 구한다.

    내용 비율 — 흰색이 아닌 픽셀의 비율. 낮으면 테두리만 있는 빈 상자다.
    흑백 여부 — 색이 없으면 투명도 마스크일 가능성이 높다.

    둘을 같이 봐야 한다. 비율만 보면 색이 진한 사진까지 걸러지고,
    흑백 여부만 보면 진짜 흑백 도표까지 걸러진다.
    """
    try:
        pix = fitz.Pixmap(image_bytes)
        is_gray = pix.n - (1 if pix.alpha else 0) == 1

        # 알파 채널이나 CMYK는 RGB로 변환한다.
        if pix.alpha or pix.colorspace is None or pix.n > 3:
            pix = fitz.Pixmap(fitz.csRGB, pix)

        # 계산을 빠르게 하기 위해 작게 줄인다.
        while pix.width > 80 and pix.height > 80:
            pix.shrink(1)

        data = pix.samples
        n = pix.n
        total = pix.width * pix.height
        if total == 0:
            return 0.0, is_gray

        ink = 0
        for i in range(0, len(data), n):
            if min(data[i], data[i + 1], data[i + 2]) < 235:
                ink += 1

        return ink / total, is_gray

    except Exception:
        # 판단할 수 없으면 일단 쓸 만하다고 본다.
        return 1.0, False


def _ink_ratio(image_bytes):
    return _image_stats(image_bytes)[0]


def _figure_rect(page, seed_side=60, include_side=18):
    """페이지에서 '그림이 있는 영역'을 하나의 사각형으로 잡는다.

    도표는 보통 사진·아이콘·화살표가 모여 하나의 뜻을 이룬다.
    그중 이미지 한 개만 떼어내면 의미가 사라지므로, 흩어진 그림들을
    묶어 그 영역을 통째로 잘라낸다.

    기준을 둘로 나눈 이유
      seed_side    이만한 그림이 하나라도 있어야 '도표가 있는 페이지'로 본다.
                   작은 아이콘만 흩어져 있는 페이지까지 잘라내면 곤란하다.
      include_side 영역을 넓힐 때는 작은 아이콘·화살표도 포함한다.
                   그렇지 않으면 도표 가장자리가 잘려 나간다.

    반환: fitz.Rect 또는 None
    """
    rects = []
    has_seed = False

    def consider(r):
        nonlocal has_seed
        if r is None:
            return
        # 머리말·꼬리말 띠는 도표가 아니다.
        if r.width > page.rect.width * 0.95 and r.height < page.rect.height * 0.15:
            return
        if r.width >= seed_side and r.height >= seed_side:
            has_seed = True
        if r.width >= include_side and r.height >= include_side:
            rects.append(r)

    for img in page.get_images(full=True):
        try:
            for r in page.get_image_rects(img[0]):
                consider(r)
        except Exception:
            continue

    # 도형으로 그린 화살표·상자도 도표의 일부다.
    try:
        for d in page.get_drawings():
            consider(d.get("rect"))
    except Exception:
        pass

    if not rects or not has_seed:
        return None

    area = rects[0]
    for r in rects[1:]:
        area = area | r          # 합집합

    # 잘린 느낌이 나지 않도록 여백을 조금 준다.
    pad = 10
    area = fitz.Rect(max(area.x0 - pad, page.rect.x0),
                     max(area.y0 - pad, page.rect.y0),
                     min(area.x1 + pad, page.rect.x1),
                     min(area.y1 + pad, page.rect.y1))

    # 너무 좁거나 페이지 전체와 다름없으면 굳이 잘라낼 이유가 없다.
    if area.width < 100 or area.height < 80:
        return None
    if area.width > page.rect.width * 0.98 and area.height > page.rect.height * 0.98:
        return None

    return area


def extract_page_image(pdf_bytes, page_no, min_size=180, min_ink=0.04,
                       max_ink=0.90, allow_page_render=True):
    """근거가 된 페이지에서 그림을 가져온다.

    1순위 : 그림들이 모여 있는 영역을 통째로 잘라낸 그림 (도표 전체)
    2순위 : 삽입된 이미지 중 내용이 있고 가장 큰 것
    3순위 : allow_page_render가 True일 때만, 페이지 전체를 렌더링한 썸네일

    영역을 잘라내는 방식을 먼저 쓰는 이유는, 화면에 보이는 그대로를 담기 때문이다.
    이미지를 직접 꺼내면 투명도 마스크가 섞이거나 도표의 일부만 나온다.

    반환: (이미지 바이트, 종류) 또는 (None, None)
    """
    try:
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    except Exception:
        return None, None

    try:
        if page_no < 1 or page_no > len(doc):
            return None, None

        page = doc[page_no - 1]

        # [1] 그림 영역을 통째로 잘라낸다. 화면에 보이는 그대로가 담긴다.
        area = _figure_rect(page)
        if area is not None:
            pix = page.get_pixmap(clip=area, matrix=fitz.Matrix(2, 2))
            png = pix.tobytes("png")
            ratio, _gray = _image_stats(png)
            if ratio >= min_ink:
                return png, "그림"

        # [2] 영역을 잡지 못하면 삽입된 이미지에서 찾는다.
        images = page.get_images(full=True)

        # 다른 이미지의 투명도 마스크로 쓰이는 것들은 그림이 아니다.
        # 목록의 두 번째 값이 그 이미지가 참조하는 마스크의 번호다.
        mask_refs = {img[1] for img in images if len(img) > 1 and img[1]}

        best = None
        for img in images:
            xref = img[0]

            if xref in mask_refs:
                continue

            try:
                info = doc.extract_image(xref)
            except Exception:
                continue

            # 로고나 아이콘처럼 작은 이미지는 제외한다.
            if info.get("width", 0) < min_size or info.get("height", 0) < min_size:
                continue

            # 흑백 1비트 이미지는 대부분 마스크다.
            if info.get("bpc", 8) == 1:
                continue

            ratio, is_gray = _image_stats(info["image"])

            # 테두리만 있는 빈 상자
            if ratio < min_ink:
                continue

            # 색이 없으면서 화면이 거의 다 채워진 것은 투명도 마스크다.
            # 컬러 사진은 흰색 아닌 픽셀이 대부분이므로 이 조건에 걸리지 않는다.
            if is_gray and ratio > max_ink:
                continue

            area = info["width"] * info["height"]
            if best is None or area > best[0]:
                best = (area, info["image"])

        if best:
            return best[1], "그림"

        if not allow_page_render:
            return None, None

        # [2] 쓸 만한 삽입 이미지가 없으면 페이지 자체를 이미지로 렌더링
        pix = page.get_pixmap(matrix=fitz.Matrix(1.6, 1.6))
        png = pix.tobytes("png")

        # 페이지가 거의 비어 있다면 그림을 넣지 않는다.
        if _ink_ratio(png) < 0.01:
            return None, None

        return png, "페이지"

    finally:
        doc.close()


# ============================================================
# 4. pptx 조립 — 공통 도구
# ============================================================
def _text(slide, left, top, width, height, content,
          size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, spacing=1.0, font=FONT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0

    p = frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = content
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return frame


def _rect(slide, left, top, width, height, fill,
          line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill

    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)

    # 기본 테마에는 그림자가 켜져 있어 지저분해 보인다. 테마 참조를 떼어내 그림자를 없앤다.
    shp.shadow.inherit = False
    style = shp._element.find(qn("p:style"))
    if style is not None:
        shp._element.remove(style)

    # 모서리 둥글기를 완만하게 (기본값은 지나치게 둥글다)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass

    return shp


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, index, title, lead):
    """모든 본문 슬라이드 상단에 공통으로 들어가는 제목 영역"""
    _text(slide, 0.62, 0.42, 3, 0.28, f"{index:02d}", size=11, bold=True, color=AMBER)
    _text(slide, 0.62, 0.72, 11.5, 0.62, title, size=25, bold=True, color=NAVY)

    if lead:
        bar = _rect(slide, 0.62, 1.5, 12.1, 0.62, ICE_LIGHT, shape=MSO_SHAPE.RECTANGLE)
        bar.line.fill.background()
        _rect(slide, 0.62, 1.5, 0.05, 0.62, AMBER, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.85, 1.5, 11.7, 0.62, lead, size=13, bold=True,
              color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# ============================================================
# 5. 슬라이드 종류별 레이아웃
# ============================================================
def _title_slide(prs, topic, subtitle):
    slide = _blank(prs)

    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK, shape=MSO_SHAPE.RECTANGLE)
    _rect(slide, 9.6, -1.5, 5.6, 5.6, NAVY_MID, shape=MSO_SHAPE.OVAL)
    _rect(slide, 11.2, 4.6, 3.2, 3.2, NAVY, shape=MSO_SHAPE.OVAL)

    _text(slide, 0.85, 2.25, 10, 0.32, "OO UNIVERSITY · 강의자료",
          size=12, bold=True, color=AMBER)
    _text(slide, 0.85, 2.7, 10.5, 1.1, topic, size=36, bold=True, color=WHITE)
    _text(slide, 0.85, 3.9, 10.5, 0.5, subtitle, size=15, color=ICE)

    _rect(slide, 0.85, 4.7, 1.6, 0.045, AMBER, shape=MSO_SHAPE.RECTANGLE)


def _agenda_slide(prs, slides):
    """목차 — 본문이 3장 이상일 때만 넣는다."""
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, LIGHT, shape=MSO_SHAPE.RECTANGLE)

    _text(slide, 0.62, 0.42, 3, 0.28, "CONTENTS", size=11, bold=True, color=AMBER)
    _text(slide, 0.62, 0.72, 11.5, 0.62, "목차", size=25, bold=True, color=NAVY)

    top = 1.75
    for i, s in enumerate(slides, start=1):
        _rect(slide, 0.62, top, 12.1, 0.78, WHITE, line=LINE)
        _rect(slide, 0.62, top, 0.06, 0.78, NAVY, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.95, top, 0.7, 0.78, f"{i:02d}", size=15, bold=True,
              color=ICE, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, 1.7, top, 10.7, 0.78, s["title"], size=14, bold=True,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
        top += 0.92


def _image_slide(prs, index, data, image_bytes, kind):
    """왼쪽 본문 + 오른쪽 그림 배치"""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    # --- 왼쪽 본문 (그림 높이에 맞춰 세로 가운데로 모은다) ---
    bullets = data["bullets"][:4]
    card_h, gap = 0.92, 0.13
    total = len(bullets) * card_h + max(len(bullets) - 1, 0) * gap
    top = 2.45 + (3.9 - total) / 2

    for bullet in bullets:
        _rect(slide, 0.62, top, 6.5, card_h, WHITE, line=LINE)
        _rect(slide, 0.62, top, 0.055, card_h, AMBER, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.92, top + 0.06, 6.0, card_h - 0.12, bullet, size=12,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)
        top += card_h + gap

    # --- 오른쪽 그림 ---
    frame_left, frame_top, frame_w, frame_h = 7.35, 2.45, 5.35, 3.9
    _rect(slide, frame_left, frame_top, frame_w, frame_h, WHITE, line=LINE)

    try:
        stream = io.BytesIO(image_bytes)
        pic = slide.shapes.add_picture(stream, Inches(frame_left + 0.18),
                                       Inches(frame_top + 0.18),
                                       width=Inches(frame_w - 0.36))

        # 세로가 넘치면 높이 기준으로 다시 맞춘다.
        max_h = Inches(frame_h - 0.36)
        if pic.height > max_h:
            ratio = pic.width / pic.height
            pic.height = int(max_h)
            pic.width = int(max_h * ratio)

        # 액자 안에서 가운데 정렬
        pic.left = Inches(frame_left) + (Inches(frame_w) - pic.width) // 2
        pic.top = Inches(frame_top) + (Inches(frame_h) - pic.height) // 2

    except Exception:
        _text(slide, frame_left, frame_top + 1.7, frame_w, 0.4,
              "이미지를 불러오지 못했습니다", size=11, color=GRAY,
              align=PP_ALIGN.CENTER)

    label = {
        "그림": "자료에서 가져온 도표",
        "페이지": "근거 페이지 미리보기",
        "삽화": "AI 생성",
    }.get(kind, "")
    _text(slide, frame_left, frame_top + frame_h + 0.12, frame_w, 0.3,
          label, size=9, color=GRAY, align=PP_ALIGN.CENTER)

    _footer(slide, data["source"])


# ============================================================
# 도식 그리기
# ============================================================
# 그림을 픽셀로 만들어내지 않고, 관계만 받아 도형으로 배치한다.
# 화살표 개수와 방향이 지시대로 정확히 그려지고, 한글도 깨지지 않는다.
# 이미지가 아니라 파워포인트 도형이므로 받는 사람이 직접 고칠 수도 있다.
def _node_box(slide, x, y, w, h, label, accent=False):
    shape = _rect(slide, x, y, w, h,
                  ICE_LIGHT if accent else WHITE,
                  line=AMBER if accent else NAVY)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.05)

    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(10.5 if len(label) > 7 else 11.5)
    run.font.bold = True
    run.font.color.rgb = NAVY
    return shape


def _line(slide, x1, y1, x2, y2, color=None, width=1.5):
    """두 점을 잇는 선."""
    line = slide.shapes.add_connector(2, Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    line.line.color.rgb = color or NAVY
    line.line.width = Pt(width)
    return line


def _arrow_head(slide, x, y, size=0.13, direction="right"):
    """연결선에 화살표 머리를 얹는다."""
    shape_id = {"right": 33, "down": 36, "left": 34, "up": 35}[direction]
    head = slide.shapes.add_shape(shape_id, Inches(x), Inches(y),
                                  Inches(size), Inches(size))
    head.fill.solid()
    head.fill.fore_color.rgb = NAVY
    head.line.fill.background()
    head.shadow.inherit = False
    style = head._element.find(qn("p:style"))
    if style is not None:
        head._element.remove(style)


def _draw_flow(slide, diagram, box):
    """흐름 — 노드를 세로로 쌓고 아래로 잇는다.

    가로로 늘어놓으면 한글 이름이 들어갈 자리가 부족해 세로로 쌓는다.
    이웃하지 않은 노드끼리 이어질 때는 오른쪽 통로로 우회시켜
    바로 아래로 가는 화살표와 겹치지 않게 한다.
    """
    left, top, width, height = box
    nodes = diagram["nodes"]

    # 우회하는 화살표가 지날 통로를 오른쪽에 비워 둔다.
    detour = any(b != a + 1 for a, b, _ in diagram["edges"])
    lane = 0.55 if detour else 0.0

    node_w = min((width - lane) * 0.86, 3.1)
    x = left + (width - lane - node_w) / 2

    node_h = 0.6
    gap = (height - node_h * len(nodes)) / max(len(nodes) - 1, 1)
    gap = max(min(gap, 0.8), 0.45)

    total = node_h * len(nodes) + gap * (len(nodes) - 1)
    y = top + (height - total) / 2

    boxes = []
    for i, name in enumerate(nodes):
        _node_box(slide, x, y, node_w, node_h, name, accent=(i == len(nodes) - 1))
        boxes.append({"cx": x + node_w / 2, "cy": y + node_h / 2,
                      "top": y, "bottom": y + node_h, "right": x + node_w})
        y += node_h + gap

    lane_x = x + node_w + lane * 0.55

    for a, b, label in diagram["edges"]:
        pa, pb = boxes[a], boxes[b]

        if b == a + 1:
            # 바로 아래로 내려가는 화살표
            _line(slide, pa["cx"], pa["bottom"], pb["cx"], pb["top"] - 0.12)
            _arrow_head(slide, pb["cx"] - 0.065, pb["top"] - 0.13, direction="down")

            if label:
                # 화살표 오른쪽에 두어 노드 글자와 겹치지 않게 한다.
                _text(slide, pa["cx"] + 0.1, (pa["bottom"] + pb["top"]) / 2 - 0.11,
                      1.1, 0.22, label, size=9, color=GRAY)
        else:
            # 오른쪽 통로로 우회한다.
            _line(slide, pa["right"], pa["cy"], lane_x, pa["cy"], AMBER, 1.4)
            _line(slide, lane_x, pa["cy"], lane_x, pb["cy"], AMBER, 1.4)
            _line(slide, lane_x, pb["cy"], pb["right"] + 0.14, pb["cy"], AMBER, 1.4)
            _arrow_head(slide, pb["right"], pb["cy"] - 0.065, direction="left")

            if label:
                _text(slide, lane_x + 0.06, (pa["cy"] + pb["cy"]) / 2 - 0.11,
                      0.7, 0.22, label, size=9, color=GRAY)


def _draw_cycle(slide, diagram, box):
    """순환 — 원 둘레에 배치하고 이웃끼리 잇는다."""
    import math

    left, top, width, height = box
    nodes = diagram["nodes"]
    n = len(nodes)

    node_w, node_h = min(width * 0.42, 1.9), 0.55
    cx, cy = left + width / 2, top + height / 2
    rx = (width - node_w) / 2 * 0.92
    ry = (height - node_h) / 2 * 0.92

    centers = []
    for i, name in enumerate(nodes):
        angle = -math.pi / 2 + 2 * math.pi * i / n
        x = cx + rx * math.cos(angle) - node_w / 2
        y = cy + ry * math.sin(angle) - node_h / 2
        _node_box(slide, x, y, node_w, node_h, name)
        centers.append((x + node_w / 2, y + node_h / 2))

    for i in range(n):
        j = (i + 1) % n
        x1, y1 = centers[i]
        x2, y2 = centers[j]

        # 선을 노드 중심까지 그으면 도형 밑으로 파고든다.
        # 양 끝을 조금씩 잘라 도형 바깥에서 시작하고 끝나게 한다.
        dx, dy = x2 - x1, y2 - y1
        dist = math.hypot(dx, dy) or 1
        ux, uy = dx / dist, dy / dist
        trim = 0.55

        sx, sy = x1 + ux * trim, y1 + uy * trim
        ex, ey = x2 - ux * trim, y2 - uy * trim

        _line(slide, sx, sy, ex, ey, NAVY, 1.4)

        # 도는 방향을 보여야 순환이라는 것이 드러난다.
        head = "right" if abs(ux) > abs(uy) else "down"
        if head == "right" and ux < 0:
            head = "left"
        if head == "down" and uy < 0:
            head = "up"
        _arrow_head(slide, ex - 0.065, ey - 0.065, direction=head)

    # 선이 노드를 덮지 않도록 노드를 다시 그린다.
    for i, name in enumerate(nodes):
        x, y = centers[i][0] - node_w / 2, centers[i][1] - node_h / 2
        _node_box(slide, x, y, node_w, node_h, name)


def _draw_tree(slide, diagram, box):
    """계층 — 첫 노드를 위에 두고 나머지를 아래에 나란히 둔다."""
    left, top, width, height = box
    nodes = diagram["nodes"]

    node_h = 0.6
    root_w = min(width * 0.6, 2.8)
    _node_box(slide, left + (width - root_w) / 2, top + 0.1, root_w, node_h,
              nodes[0], accent=True)

    children = nodes[1:]
    if not children:
        return

    child_w = min((width - 0.3 * (len(children) - 1)) / len(children), 2.4)
    total = child_w * len(children) + 0.3 * (len(children) - 1)
    x = left + (width - total) / 2
    y = top + 0.1 + node_h + 0.85

    root_cx = left + width / 2
    root_bottom = top + 0.1 + node_h

    for name in children:
        _node_box(slide, x, y, child_w, node_h, name)

        cx = x + child_w / 2
        _line(slide, root_cx, root_bottom, root_cx, root_bottom + 0.42, width=1.4)
        _line(slide, root_cx, root_bottom + 0.42, cx, root_bottom + 0.42, width=1.4)
        _line(slide, cx, root_bottom + 0.42, cx, y - 0.13, width=1.4)

        _arrow_head(slide, cx - 0.065, y - 0.14, direction="down")
        x += child_w + 0.3


def _draw_diagram(slide, diagram, box):
    """도식을 그린다. box는 (왼쪽, 위, 너비, 높이)."""
    if diagram["kind"] == "순환":
        _draw_cycle(slide, diagram, box)
    elif diagram["kind"] == "계층":
        _draw_tree(slide, diagram, box)
    else:
        _draw_flow(slide, diagram, box)


def _diagram_slide(prs, index, data):
    """왼쪽 본문 + 오른쪽 도식."""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    bullets = data["bullets"][:4]
    card_h, gap = 0.92, 0.13
    total = len(bullets) * card_h + max(len(bullets) - 1, 0) * gap
    top = 2.45 + (3.9 - total) / 2

    for bullet in bullets:
        _rect(slide, 0.62, top, 6.1, card_h, WHITE, line=LINE)
        _rect(slide, 0.62, top, 0.055, card_h, AMBER, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.92, top + 0.06, 5.6, card_h - 0.12, bullet, size=12,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)
        top += card_h + gap

    _rect(slide, 7.0, 2.35, 5.75, 4.1, LIGHT, line=LINE)
    _draw_diagram(slide, data["diagram"], (7.15, 2.5, 5.45, 3.8))

    _footer(slide, data["source"])


def _code_slide(prs, index, data):
    """코드형 — 왼쪽 설명, 오른쪽 코드 블록."""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    lines = data["code"].splitlines()[:14]
    bullets = data["bullets"][:3]

    # 설명이 없으면 코드를 넓게 쓴다.
    code_x, code_w = (6.9, 5.85) if bullets else (0.62, 12.1)

    top = 2.45
    for bullet in bullets:
        _rect(slide, 0.62, top, 6.0, 0.92, WHITE, line=LINE)
        _rect(slide, 0.62, top, 0.055, 0.92, AMBER, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.92, top + 0.06, 5.5, 0.8, bullet, size=12,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)
        top += 1.05

    # 코드 블록. 어두운 바탕에 고정폭 글꼴을 쓴다.
    height = min(4.1, 0.34 + 0.235 * max(len(lines), 1))
    _rect(slide, code_x, 2.45, code_w, height, NAVY_DARK)

    y = 2.58
    for line in lines:
        # 코드는 고정폭 글꼴이어야 들여쓰기가 어긋나지 않는다.
        _text(slide, code_x + 0.22, y, code_w - 0.44, 0.24,
              line if line.strip() else " ", size=10.5, color=ICE, font=MONO)
        y += 0.235

    _footer(slide, data["source"])


def _compare_slide(prs, index, data):
    """비교형 — 좌우 두 열로 나눠 마주 보게 둔다."""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    pairs = []
    for bullet in data["bullets"][:4]:
        left, _, right = bullet.partition("|")
        pairs.append((left.strip(), right.strip()))

    top = 2.5
    for i, (left, right) in enumerate(pairs):
        _rect(slide, 0.62, top, 5.9, 1.0, WHITE, line=LINE)
        _rect(slide, 0.62, top, 0.055, 1.0, NAVY, shape=MSO_SHAPE.RECTANGLE)
        _text(slide, 0.92, top + 0.08, 5.4, 0.84, left, size=12.5,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

        if right:
            _rect(slide, 6.85, top, 5.9, 1.0, ICE_LIGHT, line=ICE)
            _rect(slide, 6.85, top, 0.055, 1.0, AMBER, shape=MSO_SHAPE.RECTANGLE)
            _text(slide, 7.15, top + 0.08, 5.4, 0.84, right, size=12.5,
                  color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

        if i < len(pairs) - 1:
            top += 1.12
        else:
            top += 1.0

    _footer(slide, data["source"])


def _step_slide(prs, index, data):
    """절차형 — 번호를 크게 두고 세로로 이어지게 둔다."""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    steps = data["bullets"][:4]
    if not steps:
        return

    area_top, area_bottom = 2.45, 6.45
    card_h = 0.9
    gap = 0.26
    total = len(steps) * card_h + (len(steps) - 1) * gap
    top = area_top + (area_bottom - area_top - total) / 2

    for i, step in enumerate(steps, start=1):
        _rect(slide, 1.35, top, 11.4, card_h, WHITE, line=LINE)

        # 번호를 원으로 크게 표시해 순서를 드러낸다.
        circle = _rect(slide, 0.62, top + (card_h - 0.62) / 2, 0.62, 0.62,
                       NAVY, shape=MSO_SHAPE.OVAL)
        frame = circle.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(i)
        run.font.name = FONT
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = WHITE

        _text(slide, 1.65, top + 0.08, 10.8, card_h - 0.16, step, size=13.5,
              color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2)

        # 다음 단계로 이어지는 선
        if i < len(steps):
            _rect(slide, 0.91, top + card_h + 0.06, 0.04, gap - 0.12,
                  ICE, shape=MSO_SHAPE.RECTANGLE)

        top += card_h + gap

    _footer(slide, data["source"])


def _card_slide(prs, index, data):
    """그림이 없을 때 — 본문을 카드 격자로 배치해 여백을 채운다."""
    slide = _blank(prs)
    _header(slide, index, data["title"], data["lead"])

    bullets = data["bullets"][:4]
    if not bullets:
        return

    # 본문에 쓸 수 있는 세로 구간 (머리말 아래 ~ 출처 위)
    area_top, area_bottom = 2.45, 6.6

    if len(bullets) <= 2:
        # 가로로 넓은 카드. 개수가 적으면 세로 가운데로 모아 여백을 줄인다.
        card_h, gap = 1.5, 0.3
        total = len(bullets) * card_h + (len(bullets) - 1) * gap
        top = area_top + (area_bottom - area_top - total) / 2

        for i, bullet in enumerate(bullets, start=1):
            _rect(slide, 0.62, top, 12.1, card_h, WHITE, line=LINE)
            _rect(slide, 0.62, top, 0.06, card_h, NAVY, shape=MSO_SHAPE.RECTANGLE)
            _badge(slide, 0.95, top + (card_h - 0.42) / 2, i)
            _text(slide, 1.75, top + 0.15, 10.6, card_h - 0.3, bullet, size=15,
                  color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.3)
            top += card_h + gap

    elif len(bullets) == 3:
        # 3개는 세로로 쌓는다. 2열로 하면 한 칸이 비어 균형이 깨진다.
        card_h, gap = 1.1, 0.25
        total = 3 * card_h + 2 * gap
        top = area_top + (area_bottom - area_top - total) / 2

        for i, bullet in enumerate(bullets, start=1):
            _rect(slide, 0.62, top, 12.1, card_h, WHITE, line=LINE)
            _rect(slide, 0.62, top, 0.06, card_h, AMBER if i % 2 == 0 else NAVY,
                  shape=MSO_SHAPE.RECTANGLE)
            _badge(slide, 0.95, top + (card_h - 0.42) / 2, i)
            _text(slide, 1.75, top + 0.12, 10.6, card_h - 0.24, bullet, size=14,
                  color=TEXT, anchor=MSO_ANCHOR.MIDDLE, spacing=1.3)
            top += card_h + gap

    else:
        # 4개는 2열 격자
        card_h = 1.85
        gap_y = 0.3
        row_top = area_top + (area_bottom - area_top - (2 * card_h + gap_y)) / 2
        positions = [
            (0.62, row_top), (6.85, row_top),
            (0.62, row_top + card_h + gap_y), (6.85, row_top + card_h + gap_y),
        ]
        for i, bullet in enumerate(bullets):
            x, y = positions[i]
            _rect(slide, x, y, 5.87, card_h, WHITE, line=LINE)
            _rect(slide, x, y, 0.06, card_h, AMBER if i % 2 else NAVY,
                  shape=MSO_SHAPE.RECTANGLE)
            _badge(slide, x + 0.33, y + 0.32, i + 1)
            _text(slide, x + 0.33, y + 0.92, 5.2, 0.8, bullet, size=13,
                  color=TEXT, spacing=1.3)

    _footer(slide, data["source"])


def _badge(slide, left, top, number):
    circle = _rect(slide, left, top, 0.42, 0.42, NAVY, shape=MSO_SHAPE.OVAL)
    frame = circle.text_frame
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE


def _footer(slide, source):
    if not source:
        return
    _rect(slide, 0.62, 6.78, 12.1, 0.02, LINE, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, 0.62, 6.9, 12.1, 0.3, f"출처 · {source}", size=9, color=GRAY)


def _closing_slide(prs, topic):
    slide = _blank(prs)
    _rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY_DARK, shape=MSO_SHAPE.RECTANGLE)
    _rect(slide, 9.9, -1.4, 5.4, 5.4, NAVY_MID, shape=MSO_SHAPE.OVAL)

    _text(slide, 0.85, 3.0, 10, 0.7, "정리", size=30, bold=True, color=WHITE)
    _text(slide, 0.85, 3.8, 10, 0.5,
          f"본 자료는 학내 등록 자료를 근거로 자동 생성된 {topic} 강의 초안입니다.",
          size=13, color=ICE)
    _text(slide, 0.85, 4.25, 10, 0.5,
          "각 슬라이드 하단의 출처에서 원문을 확인하실 수 있습니다.",
          size=13, color=ICE)


# ============================================================
# 6. 전체 조립
# ============================================================
# 그림을 무엇으로 채울지 정하는 방식
IMAGE_MODE_DIAGRAM = "자료 도표 우선"      # 기본. 자료에 있는 그림이 곧 근거다
IMAGE_MODE_ILLUST = "AI 삽화 우선"         # 도표가 있어도 삽화를 만든다
IMAGE_MODE_NONE = "자료 그림만"            # 삽화를 만들지 않는다

IMAGE_MODES = (IMAGE_MODE_DIAGRAM, IMAGE_MODE_ILLUST, IMAGE_MODE_NONE)


def build_pptx(topic, slides, pdf_store=None, log=None,
               image_mode=IMAGE_MODE_NONE, progress=None):
    """슬라이드 목록을 pptx 파일(바이트)로 만들어 반환한다.

    pdf_store  : {파일명: PDF 바이트} — 근거 페이지의 그림을 가져오는 데 사용한다.
    log        : 슬라이드별 처리 결과를 담을 리스트 (선택).
    image_mode : 그림을 무엇으로 채울지. IMAGE_MODES 중 하나.
    progress   : 진행 상황을 알리는 함수 (선택). progress(현재, 전체, 설명)

    기본값이 '자료 도표 우선'인 이유는 자료에 있는 그림이 곧 근거이기 때문이다.
    AI 삽화는 근거가 아니라 장식이므로 슬라이드에 'AI 생성'이라고 표시한다.
    어떤 방식이든 마지막에는 근거 페이지 썸네일, 그다음 카드형으로 물러난다.
    """
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    _title_slide(prs, topic, "학내 자료 기반 강의 슬라이드 · AI 생성 초안")

    if len(slides) >= 3:
        _agenda_slide(prs, slides)

    total = len(slides)

    # 이전 작업에서 막혔던 기록을 지우고 시작한다.
    if image_mode != IMAGE_MODE_NONE:
        try:
            import imagegen
            imagegen.reset_block()
        except Exception:
            pass

    for i, data in enumerate(slides, start=1):
        image_bytes, kind = None, None
        note = ""

        file_name, page_no = _parse_source(data.get("source", ""))
        has_pdf = bool(pdf_store and file_name and page_no and file_name in pdf_store)

        def make_illustration():
            """AI 삽화를 만든다. (이미지, 종류, 설명)"""
            if progress:
                progress(i, total, f"{i}번째 슬라이드 삽화 생성 중")
            try:
                from imagegen import generate_illustration
                img, error = generate_illustration(data["title"], data["bullets"])
                if img:
                    return img, "삽화", "AI 삽화 생성"
                return None, None, f"삽화 실패 — {error}"
            except Exception as e:
                return None, None, f"삽화 실패 — {type(e).__name__}"

        def find_diagram():
            """자료에 실제로 삽입된 도표를 찾는다. 이것이 근거다."""
            if not has_pdf:
                return None, None, ""
            img, k = extract_page_image(pdf_store[file_name], page_no,
                                        allow_page_render=False)
            return (img, k, "자료 도표 사용") if img else (None, None, "")

        # 도식이 있으면 삽화를 만들지 않는다.
        # 도식은 자료의 관계를 그대로 옮긴 것이라 정확하고 비용도 들지 않는다.
        can_illustrate = not data.get("diagram")

        if image_mode == IMAGE_MODE_ILLUST and can_illustrate:
            steps = [make_illustration, find_diagram]
        elif image_mode == IMAGE_MODE_DIAGRAM and can_illustrate:
            steps = [find_diagram, make_illustration]
        else:
            steps = [find_diagram]

        for step in steps:
            image_bytes, kind, message = step()
            if message:
                note = (note + " / " + message).strip(" /")
            if image_bytes:
                break

        # [3] 그래도 없으면 근거 페이지 썸네일을 쓴다.
        if image_bytes is None and has_pdf:
            image_bytes, kind = extract_page_image(
                pdf_store[file_name], page_no, allow_page_render=True
            )
            if image_bytes:
                note = (note + " / 페이지 썸네일로 대체").strip(" /")

        if data.get("diagram") and not image_bytes:
            note = f"도식 생성 ({data['diagram']['kind']})"

        if log is not None:
            log.append(f"{i}. {data['title']} — {note or '그림 없음 (카드형)'}")

        # 유형별 배치가 내용을 드러내는 경우에는 그림보다 그쪽을 앞세운다.
        # 비교를 좌우로 나누지 않으면 "A | B" 같은 구분 기호가 그대로 보이고,
        # 절차를 번호 없이 늘어놓으면 순서라는 사실이 사라진다.
        kind_of = data.get("type")

        if kind_of == "코드" and data.get("code"):
            _code_slide(prs, i, data)
        elif kind_of == "비교":
            _compare_slide(prs, i, data)
        elif kind_of == "절차" and not data.get("diagram"):
            _step_slide(prs, i, data)
        elif image_bytes:
            _image_slide(prs, i, data, image_bytes, kind)
        elif data.get("diagram"):
            _diagram_slide(prs, i, data)
        else:
            _card_slide(prs, i, data)

    _closing_slide(prs, topic)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
