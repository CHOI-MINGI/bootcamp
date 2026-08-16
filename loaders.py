"""문서 포맷별 텍스트 추출.

지원: PDF · PPTX · DOCX · TXT · MD · HWP · HWPX

출처 표기를 유지하기 위해 포맷마다 '위치'를 페이지 번호처럼 다룬다.
    PDF        페이지 번호
    PPTX       슬라이드 번호
    그 외      문단 묶음 번호 (문서에 페이지 개념이 없음)

스캔본 PDF는 텍스트 층이 없어 추출 결과가 비어 있다.
이 경우 페이지를 이미지로 만들어 Gemini의 이미지 이해 기능으로 글자를 읽는다.
별도 OCR 엔진을 설치하지 않아 컨테이너가 가벼워진다.
"""

import base64
import os

from langchain_core.documents import Document

SUPPORTED = ["pdf", "pptx", "docx", "txt", "md", "hwp", "hwpx"]

# DOCX·TXT는 페이지가 없어 일정 분량마다 위치 번호를 올린다.
CHARS_PER_UNIT = 1500

# 이 글자 수보다 적으면 스캔 이미지 페이지로 보고 OCR을 시도한다.
OCR_THRESHOLD = 30


def ocr_enabled():
    return os.getenv("OCR_ENABLED", "1").strip() not in ("0", "false", "False")


def _ocr(page):
    """텍스트가 없는 페이지를 이미지로 만들어 글자를 읽어낸다.

    별도 OCR 엔진을 설치하지 않고 이미 쓰고 있는 Gemini의 이미지 이해 기능을 쓴다.
    설치할 시스템 패키지가 없어 컨테이너가 가벼워진다.
    """
    import fitz
    from langchain_core.messages import HumanMessage
    from langchain_google_genai import ChatGoogleGenerativeAI

    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
    data = base64.b64encode(pix.tobytes("png")).decode()

    llm = ChatGoogleGenerativeAI(
        model=os.getenv("CHAT_MODEL", "gemini-3.6-flash"), temperature=0
    )

    message = HumanMessage(content=[
        {"type": "text",
         "text": "이 이미지에 있는 글자를 그대로 옮겨 적으십시오. "
                 "설명이나 요약을 덧붙이지 말고 본문만 출력하십시오. "
                 "글자가 없으면 아무것도 출력하지 마십시오."},
        {"type": "image_url", "image_url": f"data:image/png;base64,{data}"},
    ])

    return llm.invoke([message]).content.strip()


def _pdf(path):
    import fitz

    docs = []
    ocr_pages = []

    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf):
            text = page.get_text().strip()

            # 스캔본은 텍스트 층이 없어 추출 결과가 비어 있다.
            if len(text) < OCR_THRESHOLD and ocr_enabled():
                try:
                    read = _ocr(page)
                    if len(read) > len(text):
                        text = read
                        ocr_pages.append(i + 1)
                except Exception:
                    # OCR에 실패해도 나머지 페이지 처리는 계속한다.
                    pass

            if text:
                meta = {"page": i}
                if (i + 1) in ocr_pages:
                    meta["ocr"] = True
                docs.append(Document(page_content=text, metadata=meta))

    return docs


def _pptx(path):
    from pptx import Presentation

    docs = []
    prs = Presentation(path)

    for i, slide in enumerate(prs.slides):
        parts = []

        for shape in slide.shapes:
            # 도형 안의 글자
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)

            # 표는 셀을 줄 단위로 붙인다.
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))

        # 발표자 노트도 강의 내용이므로 함께 담는다.
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                parts.append(f"[발표자 노트] {note}")

        if parts:
            docs.append(Document(page_content="\n".join(parts),
                                 metadata={"page": i}))
    return docs


def _docx(path):
    import docx

    document = docx.Document(path)
    blocks = [p.text.strip() for p in document.paragraphs if p.text.strip()]

    # 표 내용도 포함한다.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                blocks.append(" | ".join(cells))

    return _group(blocks)


def _txt(path):
    for encoding in ("utf-8", "utf-8-sig", "cp949"):
        try:
            with open(path, encoding=encoding) as f:
                blocks = [line.strip() for line in f if line.strip()]
            return _group(blocks)
        except UnicodeDecodeError:
            continue

    raise ValueError("텍스트 인코딩을 인식하지 못했습니다.")


def _hwp(path):
    """한글 문서(HWP 5.0). OLE 복합 문서 안의 본문 스트림을 직접 읽는다."""
    import struct
    import zlib

    import olefile

    ole = olefile.OleFileIO(path)

    try:
        # 파일 헤더의 플래그로 본문 압축 여부를 알 수 있다.
        header = ole.openstream("FileHeader").read()
        compressed = bool(header[36] & 1)

        sections = sorted(
            (e for e in ole.listdir() if e[0] == "BodyText"),
            key=lambda e: int(e[1].replace("Section", "")),
        )

        blocks = []

        for entry in sections:
            data = ole.openstream(entry).read()
            if compressed:
                data = zlib.decompress(data, -15)

            blocks.extend(_hwp_records(data, struct))

        if not blocks:
            raise ValueError("본문을 찾지 못했습니다.")

        return _group(blocks)

    finally:
        ole.close()


def _hwp_records(data, struct):
    """HWP 본문 스트림을 레코드 단위로 훑어 문단 텍스트만 모은다."""
    PARA_TEXT = 67          # 문단 본문을 담는 레코드 종류
    blocks = []
    cursor = 0

    while cursor < len(data):
        head = struct.unpack_from("<I", data, cursor)[0]
        tag = head & 0x3FF
        size = (head >> 20) & 0xFFF
        cursor += 4

        # 크기가 4095를 넘으면 다음 4바이트에 실제 크기가 들어 있다.
        if size == 0xFFF:
            size = struct.unpack_from("<I", data, cursor)[0]
            cursor += 4

        if tag == PARA_TEXT:
            raw = data[cursor:cursor + size].decode("utf-16le", errors="ignore")
            # 서식·표 위치를 나타내는 제어 문자를 걸러낸다.
            text = "".join(c for c in raw if ord(c) > 31 or c in "\n\t").strip()
            if text:
                blocks.append(text)

        cursor += size

    return blocks


def _hwpx(path):
    """한글 문서(HWPX). 압축된 XML 묶음이라 훨씬 단순하다."""
    import re as _re
    import zipfile

    blocks = []

    with zipfile.ZipFile(path) as z:
        names = sorted(n for n in z.namelist()
                       if n.startswith("Contents/section") and n.endswith(".xml"))

        for name in names:
            xml = z.read(name).decode("utf-8", errors="ignore")

            # <hp:t> 태그 안에 실제 글자가 들어 있다.
            for match in _re.findall(r"<hp:t[^>]*>(.*?)</hp:t>", xml, _re.S):
                text = _re.sub(r"<[^>]+>", "", match).strip()
                if text:
                    blocks.append(text)

    if not blocks:
        raise ValueError("본문을 찾지 못했습니다.")

    return _group(blocks)


def _group(blocks):
    """페이지가 없는 문서를 일정 분량 단위로 묶어 위치 번호를 붙인다."""
    docs = []
    buffer = []
    length = 0

    def flush():
        nonlocal buffer, length
        if buffer:
            docs.append(Document(page_content="\n".join(buffer),
                                 metadata={"page": len(docs)}))
            buffer, length = [], 0

    for block in blocks:
        # 줄바꿈 없이 아주 길게 이어진 문단은 그대로 두면 통째로 한 구간이 된다.
        # 그러면 출처가 "구간 1" 하나로 뭉뚱그려지므로 미리 잘라 준다.
        while len(block) > CHARS_PER_UNIT:
            flush()
            docs.append(Document(page_content=block[:CHARS_PER_UNIT],
                                 metadata={"page": len(docs)}))
            block = block[CHARS_PER_UNIT:]

        buffer.append(block)
        length += len(block)

        if length >= CHARS_PER_UNIT:
            flush()

    flush()
    return docs


LOADERS = {
    "pdf": _pdf,
    "pptx": _pptx,
    "docx": _docx,
    "txt": _txt,
    "md": _txt,
    "hwp": _hwp,
    "hwpx": _hwpx,
}


def extension(file_name):
    return os.path.splitext(file_name)[1].lower().lstrip(".")


def load(path, file_name):
    """파일에서 텍스트를 추출해 Document 목록으로 돌려준다."""
    ext = extension(file_name)

    loader = LOADERS.get(ext)
    if loader is None:
        raise ValueError(f"지원하지 않는 형식입니다 — .{ext}")

    docs = loader(path)

    if not docs:
        raise ValueError("문서에서 텍스트를 찾지 못했습니다. "
                         "빈 파일이거나 형식이 손상되었을 수 있습니다.")

    return docs


def location_label(file_name):
    """출처에 쓸 위치 표기. 포맷마다 부르는 이름이 다르다."""
    return {
        "pdf": "p.",
        "pptx": "슬라이드 ",
    }.get(extension(file_name), "구간 ")
